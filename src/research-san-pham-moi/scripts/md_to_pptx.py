"""
Convert .md research report -> .pptx slides.
Author: Techla — v1.0.0 — License: xem LICENSE.md

Chiến lược:
1. Nếu có `marp` CLI trong PATH -> dùng Marp (đẹp hơn, theme pro).
2. Nếu không -> fallback python-pptx (đơn giản, không phụ thuộc Node).

Cách chia slide (fallback):
- H1 hoặc H2 -> slide mới (title = heading text)
- H3 -> section heading trong slide (vẫn mở rộng slide hiện tại)
- Body text / bullets -> content slide

Usage:
    python scripts/md_to_pptx.py --md research.md --pptx research.pptx

Flags:
    --engine marp|pptx|auto   (default auto — dùng marp nếu có)
    --theme default|gaia|uncover (chỉ áp dụng khi engine=marp)
"""
from __future__ import annotations
import os
import sys
import re
import shutil
import subprocess
import argparse
from pathlib import Path


FOOTER_TEXT = "Techla — Skill Research sản phẩm mới v1.0.0"


# ═══════════════════════════════════════════════════════════════════
# ENGINE 1: Marp CLI
# ═══════════════════════════════════════════════════════════════════

def marp_available() -> bool:
    return shutil.which("marp") is not None


def _wrap_marp_md(md_text: str, theme: str = "default") -> str:
    """Thêm front-matter Marp + chia slide tự động dựa trên heading."""
    front_matter = (
        "---\n"
        "marp: true\n"
        f"theme: {theme}\n"
        "paginate: true\n"
        f"footer: '{FOOTER_TEXT}'\n"
        "size: 16:9\n"
        "---\n\n"
    )

    # Nếu chưa có "---" phân chia slide, tự chèn trước mỗi H2
    lines = md_text.splitlines()
    out = []
    first_h2_seen = False
    for line in lines:
        if line.startswith("## "):
            if first_h2_seen:
                out.append("\n---\n")
            first_h2_seen = True
        out.append(line)
    body = "\n".join(out)
    return front_matter + body


def build_with_marp(md_path: Path, pptx_path: Path, theme: str = "default") -> bool:
    md_text = md_path.read_text(encoding="utf-8")
    wrapped = _wrap_marp_md(md_text, theme=theme)

    tmp_md = pptx_path.with_suffix(".marp.md")
    tmp_md.write_text(wrapped, encoding="utf-8")

    try:
        result = subprocess.run(
            ["marp", str(tmp_md), "--pptx", "-o", str(pptx_path), "--allow-local-files"],
            capture_output=True, text=True, timeout=120,
        )
        if result.returncode != 0:
            print(f"Marp thất bại: {result.stderr}", file=sys.stderr)
            return False
        return True
    except Exception as e:
        print(f"Lỗi khi chạy Marp: {e}", file=sys.stderr)
        return False
    finally:
        if tmp_md.exists():
            try:
                tmp_md.unlink()
            except OSError:
                pass


# ═══════════════════════════════════════════════════════════════════
# ENGINE 2: python-pptx fallback
# ═══════════════════════════════════════════════════════════════════

HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
BULLET_RE = re.compile(r"^\s*[-*]\s+(.+)$")
NUM_RE = re.compile(r"^\s*\d+\.\s+(.+)$")
BLOCKQUOTE_RE = re.compile(r"^\s*>\s+(.+)$")


def _strip_md(text: str) -> str:
    # Bỏ **bold**, *italic*, `code`, [link](url), [^refs]
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*(?!\s)(.+?)(?<!\s)\*(?!\*)", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[\^(\d+)\]", r"[\1]", text)
    text = re.sub(r"\[(.+?)\]\((.+?)\)", r"\1", text)
    return text.strip()


def _parse_slides(md_text: str) -> list[dict]:
    """
    Chia md thành các slide dựa trên heading level 1-2.
    Mỗi slide: {"title": str, "bullets": [str], "notes": [str]}
    """
    slides = []
    current = None

    for raw in md_text.splitlines():
        line = raw.rstrip()
        if not line.strip():
            if current and current["bullets"] and current["bullets"][-1] != "":
                pass
            continue

        m = HEADING_RE.match(line)
        if m:
            level = len(m.group(1))
            title = _strip_md(m.group(2))
            if level <= 2:
                if current:
                    slides.append(current)
                current = {"title": title, "level": level, "bullets": []}
                continue
            else:
                # H3+ -> sub-heading
                if current is None:
                    current = {"title": title, "level": level, "bullets": []}
                else:
                    current["bullets"].append(("heading", title))
                continue

        if current is None:
            current = {"title": "Mở đầu", "level": 1, "bullets": []}

        mb = BULLET_RE.match(line)
        if mb:
            current["bullets"].append(("bullet", _strip_md(mb.group(1))))
            continue
        mn = NUM_RE.match(line)
        if mn:
            current["bullets"].append(("bullet", _strip_md(mn.group(1))))
            continue
        mq = BLOCKQUOTE_RE.match(line)
        if mq:
            current["bullets"].append(("quote", _strip_md(mq.group(1))))
            continue
        if line.startswith("---") or line.startswith("___"):
            continue

        cleaned = _strip_md(line)
        if cleaned:
            current["bullets"].append(("text", cleaned))

    if current:
        slides.append(current)

    return slides


def build_with_pptx(md_path: Path, pptx_path: Path) -> bool:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("Lỗi: thiếu python-pptx. Cài: pip install python-pptx", file=sys.stderr)
        return False

    md_text = md_path.read_text(encoding="utf-8")
    slides_data = _parse_slides(md_text)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    TITLE_COLOR = RGBColor(0x1F, 0x4E, 0x79)
    BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
    QUOTE_COLOR = RGBColor(0x7F, 0x8C, 0x8D)

    blank_layout = prs.slide_layouts[6]

    # --- Slide 1: title slide ---
    first = slides_data[0] if slides_data else {"title": "Báo cáo Research", "bullets": []}
    s = prs.slides.add_slide(blank_layout)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = first["title"]
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = FOOTER_TEXT
    r2.font.size = Pt(14)
    r2.font.italic = True
    r2.font.color.rgb = QUOTE_COLOR

    # --- Slides 2+: content ---
    for sd in slides_data[1:]:
        s = prs.slides.add_slide(blank_layout)

        # Title
        title_box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = sd["title"]
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = TITLE_COLOR

        # Body
        body_box = s.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.5))
        tfb = body_box.text_frame
        tfb.word_wrap = True
        first_body_para = True

        for kind, text in sd["bullets"]:
            if not text:
                continue
            if first_body_para:
                pp = tfb.paragraphs[0]
                first_body_para = False
            else:
                pp = tfb.add_paragraph()

            rr = pp.add_run()
            if kind == "heading":
                rr.text = text
                rr.font.size = Pt(20)
                rr.font.bold = True
                rr.font.color.rgb = TITLE_COLOR
                pp.space_after = Pt(4)
            elif kind == "bullet":
                rr.text = "• " + text
                rr.font.size = Pt(18)
                rr.font.color.rgb = BODY_COLOR
                pp.space_after = Pt(3)
            elif kind == "quote":
                rr.text = f"\u201C{text}\u201D"
                rr.font.size = Pt(16)
                rr.font.italic = True
                rr.font.color.rgb = QUOTE_COLOR
                pp.space_after = Pt(6)
            else:
                rr.text = text
                rr.font.size = Pt(18)
                rr.font.color.rgb = BODY_COLOR
                pp.space_after = Pt(6)

        # Footer
        ft_box = s.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
        ftf = ft_box.text_frame
        pft = ftf.paragraphs[0]
        pft.alignment = PP_ALIGN.RIGHT
        rft = pft.add_run()
        rft.text = FOOTER_TEXT
        rft.font.size = Pt(9)
        rft.font.italic = True
        rft.font.color.rgb = QUOTE_COLOR

    prs.save(str(pptx_path))
    return True


# ═══════════════════════════════════════════════════════════════════
# CLI
# ═══════════════════════════════════════════════════════════════════

def main():
    ap = argparse.ArgumentParser(description="Convert Markdown research report -> PowerPoint (.pptx).")
    ap.add_argument("--md", required=True, help="Input .md file")
    ap.add_argument("--pptx", required=True, help="Output .pptx path")
    ap.add_argument("--engine", choices=["marp", "pptx", "auto"], default="auto")
    ap.add_argument("--theme", default="default", help="Marp theme (default|gaia|uncover)")
    args = ap.parse_args()

    md_path = Path(args.md)
    pptx_path = Path(args.pptx)
    if not md_path.exists():
        print(f"Lỗi: không tìm thấy '{md_path}'.", file=sys.stderr); sys.exit(1)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    engine = args.engine
    if engine == "auto":
        engine = "marp" if marp_available() else "pptx"

    print(f"Engine: {engine}")
    if engine == "marp":
        if not marp_available():
            print("Lỗi: marp CLI không có trong PATH. Cài bằng: npm i -g @marp-team/marp-cli",
                  file=sys.stderr); sys.exit(2)
        ok = build_with_marp(md_path, pptx_path, theme=args.theme)
    else:
        ok = build_with_pptx(md_path, pptx_path)

    if not ok:
        sys.exit(3)
    print(f"OK: {pptx_path}")


if __name__ == "__main__":
    main()
